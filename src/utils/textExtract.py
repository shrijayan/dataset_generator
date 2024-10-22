import os
import json
import xml.etree.ElementTree as ET
from PyPDF2 import PdfReader
from pptx import Presentation
from bs4 import BeautifulSoup
import shutil

with open('config.json', 'r') as config_file:
    config = json.load(config_file)

class TextExtractor():
    def __init__(self):
        folder_path = f"{config.get('chroma_db_path')}"
        if os.path.exists(folder_path):
            shutil.rmtree(folder_path)
            print(f"Folder '{folder_path}' has been deleted successfully.")
        else:
            print(f"Folder '{folder_path}' does not exist.")
        
    def extract_text_from_file(self, file_path):
        _, ext = os.path.splitext(file_path)
        ext = ext.lower()

        try:
            if ext == '.txt':
                with open(file_path, 'r', encoding='utf-8') as f:
                    return f.read()
            elif ext == '.json':
                with open(file_path, 'r', encoding='utf-8') as f:
                    data = json.load(f)
                return json.dumps(data, indent=2)
            elif ext == '.xml':
                tree = ET.parse(file_path)
                return ET.tostring(tree.getroot(), encoding='unicode', method='xml')
            elif ext == '.pdf':
                text = []
                pdf = PdfReader(file_path)
                for page in pdf.pages:
                    # Extract text and replace multiple newlines with spaces
                    page_text = page.extract_text()
                    cleaned_text = ' '.join(page_text.splitlines())  # Replace line breaks with spaces
                    text.append(cleaned_text)
                return ' '.join(text)
            elif ext in ['.ppt', '.pptx']:
                text = []
                prs = Presentation(file_path)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, 'text'):
                            text.append(shape.text)
                return ' '.join(text)
            elif ext == '.html':
                with open(file_path, 'r', encoding='utf-8') as f:
                    soup = BeautifulSoup(f, 'html.parser')
                return soup.get_text()
            else:
                return f"Unsupported file type: {ext}"
        except Exception as e:
            return f"Error processing file: {str(e)}"

    def process_folder(self, folder_path):
        for filename in os.listdir(folder_path):
            file_path = os.path.join(folder_path, filename)
            
            # Skip directories
            if os.path.isdir(file_path):
                continue
            
            # Process non-txt files
            if not filename.lower().endswith('.txt'):
                try:
                    # Extract text from the file
                    text = self.extract_text_from_file(file_path)
                    
                    # Create a new txt file with the extracted text
                    new_filename = os.path.splitext(filename)[0] + '.txt'
                    new_file_path = os.path.join(folder_path, new_filename)
                    
                    with open(new_file_path, 'w', encoding='utf-8') as f:
                        f.write(text)
                        
                except Exception as e:
                    print(f"Error processing {filename}: {str(e)}")
            else:
                print(f"Skipped: {filename} (already a .txt file)")
                pass

if __name__ == "__main__":
    folder_path = input("Enter the folder path: ")
    text_extractor = TextExtractor()
    text_extractor.process_folder(folder_path)